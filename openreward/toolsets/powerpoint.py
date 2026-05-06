from __future__ import annotations

import json
from typing import Literal, Optional
from pydantic import BaseModel, Field

from openreward.environments.environment import tool
from openreward.environments.toolset import Toolset
from openreward.environments.types import TextBlock, ToolOutput


# ===== Pydantic 参数模型 =====

class ReadAllParams(BaseModel):
    file_path: str = Field(..., description="沙箱中 PowerPoint 文件的路径（例如，/orwd_data/presentation.pptx）")


class ReadSlidesParams(BaseModel):
    file_path: str = Field(..., description="沙箱中 PowerPoint 文件的路径")
    slide_indices: Optional[list[int]] = Field(None, description="要读取的幻灯片索引列表（从 0 开始）。如果为 None，则读取所有幻灯片。")


class CreatePresentationParams(BaseModel):
    file_path: str = Field(..., description="沙箱中新演示文稿的输出路径")
    title: Optional[str] = Field(None, description="第一张幻灯片的可选标题文本")


class DeletePresentationParams(BaseModel):
    file_path: str = Field(..., description="沙箱中要删除的演示文稿文件路径")


class AddSlideParams(BaseModel):
    file_path: str = Field(..., description="沙箱中演示文稿的路径")
    layout: Literal["title", "title_content", "section_header", "two_content", "comparison", "title_only", "blank", "content_with_caption", "picture_with_caption"] = Field(..., description="幻灯片布局类型")
    title: Optional[str] = Field(None, description="幻灯片标题文本")
    content: Optional[str] = Field(None, description="幻灯片正文/内容文本")


class EditSlideParams(BaseModel):
    file_path: str = Field(..., description="沙箱中演示文稿的路径")
    slide_index: int = Field(..., description="要编辑的从零开始的幻灯片索引")
    title: Optional[str] = Field(None, description="新标题文本（更新标题形状/占位符）")
    content: Optional[str] = Field(None, description="新正文内容（更新内容占位符）")


class AddImageParams(BaseModel):
    file_path: str = Field(..., description="沙箱中演示文稿的路径")
    slide_index: int = Field(..., description="要添加图像的从零开始的幻灯片索引")
    image_path: str = Field(..., description="沙箱中图像文件的路径")
    left: float = Field(..., description="左侧位置（英寸）")
    top: float = Field(..., description="顶部位置（英寸）")
    width: Optional[float] = Field(None, description="宽度（英寸）（None 表示自动）")
    height: Optional[float] = Field(None, description="高度（英寸）（None 表示自动）")


class ModifyImageParams(BaseModel):
    file_path: str = Field(..., description="沙箱中演示文稿的路径")
    slide_index: int = Field(..., description="从零开始的幻灯片索引")
    image_index: int = Field(..., description="幻灯片上图像的从零开始的索引")
    new_image_path: Optional[str] = Field(None, description="替换图像的路径（None 表示保留现有图像）")
    width: Optional[float] = Field(None, description="新宽度（英寸）")
    height: Optional[float] = Field(None, description="新高度（英寸）")


class ChartSeriesData(BaseModel):
    name: str = Field(..., description="系列名称")
    values: list[float] = Field(..., description="系列的数据值")


class AddChartParams(BaseModel):
    file_path: str = Field(..., description="沙箱中演示文稿的路径")
    slide_index: int = Field(..., description="从零开始的幻灯片索引")
    chart_type: Literal["column", "bar", "line", "pie"] = Field(..., description="图表类型")
    categories: list[str] = Field(..., description="类别标签（x 轴）")
    series: list[ChartSeriesData] = Field(..., description="图表的数据系列")
    left: float = Field(..., description="左侧位置（英寸）")
    top: float = Field(..., description="顶部位置（英寸）")
    width: float = Field(..., description="宽度（英寸）")
    height: float = Field(..., description="高度（英寸）")


class AddTableParams(BaseModel):
    file_path: str = Field(..., description="沙箱中演示文稿的路径")
    slide_index: int = Field(..., description="从零开始的幻灯片索引")
    rows: int = Field(..., description="行数")
    cols: int = Field(..., description="列数")
    data: list[list[str]] = Field(..., description="单元格内容的二维数组 [行][列]")
    left: float = Field(..., description="左侧位置（英寸）")
    top: float = Field(..., description="顶部位置（英寸）")
    width: float = Field(..., description="宽度（英寸）")
    height: float = Field(..., description="高度（英寸）")


class ListImagesParams(BaseModel):
    file_path: str = Field(..., description="沙箱中演示文稿的路径")
    slide_index: Optional[int] = Field(None, description="特定幻灯片索引（None 表示所有幻灯片）")


# ===== PowerPoint 工具集类 =====

class PowerPointToolset(Toolset):
    """通过 python-pptx 在沙箱中执行，提供 PowerPoint 操作工具的工具集"""

    @tool
    async def powerpoint_read_all(self, params: ReadAllParams) -> ToolOutput:
        """读取完整的演示文稿结构，包括元数据和所有幻灯片内容"""
        script = f'''
import json
from pptx import Presentation

try:
    prs = Presentation("{params.file_path}")

    slides_data = []
    for idx, slide in enumerate(prs.slides):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            text_parts.append(run.text)

        slides_data.append({{
            "index": idx,
            "layout": slide.slide_layout.name,
            "text": "\\n".join(text_parts),
            "shape_count": len(slide.shapes),
        }})

    result = {{
        "success": True,
        "slide_count": len(prs.slides),
        "width": prs.slide_width / 914400,
        "height": prs.slide_height / 914400,
        "slides": slides_data,
    }}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error reading presentation: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            # 构建摘要
            summary = f"Presentation: {params.file_path}\n"
            summary += f"Slides: {result['slide_count']}\n"
            summary += f'Dimensions: {result["width"]:.2f}" x {result["height"]:.2f}"\n\n'
            summary += "Slide Overview:\n"
            for s in result['slides'][:5]:
                preview = s['text'][:60].replace("\n", " ")
                summary += f"  [{s['index']}] {s['layout']}: {preview}...\n"
            if result['slide_count'] > 5:
                summary += f"  ... and {result['slide_count'] - 5} more slides\n"

            return ToolOutput(
                blocks=[TextBlock(text=summary)],
                metadata={
                    "file_path": params.file_path,
                    "slide_count": result['slide_count'],
                    "dimensions": {"width": result['width'], "height": result['height']},
                    "slides": result['slides'],
                },
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_read_slides(self, params: ReadSlidesParams) -> ToolOutput:
        """从特定幻灯片或所有幻灯片读取文本内容"""
        indices_str = json.dumps(params.slide_indices) if params.slide_indices else "None"

        script = f'''
import json
from pptx import Presentation

try:
    prs = Presentation("{params.file_path}")
    indices = {indices_str}

    if indices is None:
        indices = list(range(len(prs.slides)))

    slides_data = []
    for idx in indices:
        if idx < 0 or idx >= len(prs.slides):
            slides_data.append({{"index": idx, "error": "Slide index out of range"}})
            continue

        slide = prs.slides[idx]
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            text_parts.append(run.text)

        slides_data.append({{
            "index": idx,
            "layout": slide.slide_layout.name,
            "text": "\\n".join(text_parts),
        }})

    result = {{"success": True, "slides": slides_data}}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            summary = f"Read {len(result['slides'])} slide(s) from {params.file_path}\n\n"
            for s in result['slides']:
                if "error" in s:
                    summary += f"Slide {s['index']}: {s['error']}\n"
                else:
                    text_preview = s['text'].replace("\n", " ")
                    summary += f"Slide {s['index']} ({s['layout']}):\n{text_preview}"

            return ToolOutput(
                blocks=[TextBlock(text=summary)],
                metadata={"file_path": params.file_path, "slides": result['slides']},
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_create(self, params: CreatePresentationParams) -> ToolOutput:
        """创建带有可选标题幻灯片的新 PowerPoint 演示文稿"""
        title_escaped = params.title.replace('\\', '\\\\').replace('"', '\\"').replace('\n', '\\n') if params.title else ""

        script = f'''
import json
from pptx import Presentation

try:
    prs = Presentation()

    if "{title_escaped}":
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = "{title_escaped}"

    prs.save("{params.file_path}")

    result = {{"success": True, "slide_count": len(prs.slides)}}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            return ToolOutput(
                blocks=[TextBlock(text=f"✅ Presentation created successfully at {params.file_path}")],
                metadata={
                    "file_path": params.file_path,
                    "slide_count": result['slide_count'],
                },
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_delete(self, params: DeletePresentationParams) -> ToolOutput:
        """从沙箱中删除 PowerPoint 演示文稿文件"""
        output, exit_code = await self.sandbox.run(f"rm {params.file_path}", max_bytes=1_000_000)

        if exit_code == 0:
            return ToolOutput(
                blocks=[TextBlock(text=f"✅ Presentation deleted successfully: {params.file_path}")],
                metadata={"file_path": params.file_path, "deleted": True},
                reward=0.0,
                finished=False,
            )
        else:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Failed to delete: {output}")],
                metadata={"file_path": params.file_path, "deleted": False, "error": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_add_slide(self, params: AddSlideParams) -> ToolOutput:
        """向演示文稿添加具有指定布局和内容的新幻灯片"""
        layout_map = {
            "title": 0,
            "title_content": 1,
            "section_header": 2,
            "two_content": 3,
            "comparison": 4,
            "title_only": 5,
            "blank": 6,
            "content_with_caption": 7,
            "picture_with_caption": 8,
        }
        layout_idx = layout_map.get(params.layout, 1)

        # 为 Python 字符串字面量转义引号和换行符
        title_escaped = params.title.replace('\\', '\\\\').replace('"', '\\"').replace('\n', '\\n') if params.title else ""
        content_escaped = params.content.replace('\\', '\\\\').replace('"', '\\"').replace('\n', '\\n') if params.content else ""

        script = f'''
import json
from pptx import Presentation

try:
    prs = Presentation("{params.file_path}")
    slide_layout = prs.slide_layouts[{layout_idx}]
    slide = prs.slides.add_slide(slide_layout)

    if "{title_escaped}" and slide.shapes.title:
        slide.shapes.title.text = "{title_escaped}"

    if "{content_escaped}":
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = "{content_escaped}"
                break

    prs.save("{params.file_path}")

    result = {{
        "success": True,
        "slide_index": len(prs.slides) - 1,
        "slide_count": len(prs.slides),
    }}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            return ToolOutput(
                blocks=[TextBlock(text=f"✅ Slide added successfully at index {result['slide_index']}")],
                metadata={
                    "file_path": params.file_path,
                    "slide_index": result['slide_index'],
                    "slide_count": result['slide_count'],
                    "layout": params.layout,
                },
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_edit_slide(self, params: EditSlideParams) -> ToolOutput:
        """编辑现有幻灯片的文本内容"""
        title_escaped = params.title.replace('\\', '\\\\').replace('"', '\\"').replace('\n', '\\n') if params.title else None
        content_escaped = params.content.replace('\\', '\\\\').replace('"', '\\"').replace('\n', '\\n') if params.content else None

        script = f'''
import json
from pptx import Presentation

try:
    prs = Presentation("{params.file_path}")

    if {params.slide_index} < 0 or {params.slide_index} >= len(prs.slides):
        raise ValueError(f"Slide index {params.slide_index} out of range")

    slide = prs.slides[{params.slide_index}]

    if "{title_escaped}" != "None" and slide.shapes.title:
        slide.shapes.title.text = "{title_escaped}"

    if "{content_escaped}" != "None":
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = "{content_escaped}"
                break

    prs.save("{params.file_path}")

    result = {{"success": True}}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            return ToolOutput(
                blocks=[TextBlock(text=f"✅ Slide {params.slide_index} updated successfully")],
                metadata={"file_path": params.file_path, "slide_index": params.slide_index},
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_add_table(self, params: AddTableParams) -> ToolOutput:
        """向幻灯片添加具有指定数据的表格"""
        data_json = json.dumps(params.data)

        script = f'''
import json
from pptx import Presentation
from pptx.util import Inches

try:
    prs = Presentation("{params.file_path}")
    slide = prs.slides[{params.slide_index}]

    data = {data_json}
    table_shape = slide.shapes.add_table(
        {params.rows}, {params.cols},
        Inches({params.left}), Inches({params.top}),
        Inches({params.width}), Inches({params.height})
    )
    table = table_shape.table

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_value in enumerate(row_data):
            table.cell(row_idx, col_idx).text = str(cell_value)

    prs.save("{params.file_path}")

    result = {{"success": True}}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            return ToolOutput(
                blocks=[TextBlock(text=f"✅ Table ({params.rows}x{params.cols}) added to slide {params.slide_index}")],
                metadata={
                    "file_path": params.file_path,
                    "slide_index": params.slide_index,
                    "rows": params.rows,
                    "cols": params.cols,
                },
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_add_image(self, params: AddImageParams) -> ToolOutput:
        """向幻灯片添加图像"""
        width_param = f"width=Inches({params.width})" if params.width else ""
        height_param = f"height=Inches({params.height})" if params.height else ""
        size_params = f", {width_param}" if width_param else ""
        size_params += f", {height_param}" if height_param else ""

        script = f'''
import json
from pptx import Presentation
from pptx.util import Inches

try:
    prs = Presentation("{params.file_path}")
    slide = prs.slides[{params.slide_index}]

    pic = slide.shapes.add_picture(
        "{params.image_path}",
        Inches({params.left}),
        Inches({params.top}){size_params}
    )

    width_inches = pic.width / 914400
    height_inches = pic.height / 914400

    prs.save("{params.file_path}")

    result = {{
        "success": True,
        "width": width_inches,
        "height": height_inches,
    }}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error"), "image_added": False},
                    reward=0.0,
                    finished=False,
                )

            return ToolOutput(
                blocks=[TextBlock(text=f"✅ Image added to slide {params.slide_index}")],
                metadata={
                    "file_path": params.file_path,
                    "slide_index": params.slide_index,
                    "image_added": True,
                    "dimensions": {"width": result['width'], "height": result['height']},
                },
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output, "image_added": False},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_add_chart(self, params: AddChartParams) -> ToolOutput:
        """向幻灯片添加图表"""
        categories_json = json.dumps(params.categories)
        series_json = json.dumps([{"name": s.name, "values": s.values} for s in params.series])

        chart_type_map = {
            "column": "XL_CHART_TYPE.COLUMN_CLUSTERED",
            "bar": "XL_CHART_TYPE.BAR_CLUSTERED",
            "line": "XL_CHART_TYPE.LINE",
            "pie": "XL_CHART_TYPE.PIE",
        }
        chart_type_code = chart_type_map.get(params.chart_type, "XL_CHART_TYPE.COLUMN_CLUSTERED")

        script = f'''
import json
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

try:
    prs = Presentation("{params.file_path}")
    slide = prs.slides[{params.slide_index}]

    chart_data = ChartData()
    chart_data.categories = {categories_json}

    series_list = {series_json}
    for series in series_list:
        chart_data.add_series(series["name"], series["values"])

    chart_shape = slide.shapes.add_chart(
        {chart_type_code},
        Inches({params.left}), Inches({params.top}),
        Inches({params.width}), Inches({params.height}),
        chart_data
    )

    prs.save("{params.file_path}")

    result = {{
        "success": True,
        "series_count": len(series_list),
        "category_count": len({categories_json}),
    }}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            return ToolOutput(
                blocks=[TextBlock(text=f"✅ {params.chart_type.capitalize()} chart added to slide {params.slide_index}")],
                metadata={
                    "file_path": params.file_path,
                    "slide_index": params.slide_index,
                    "chart_type": params.chart_type,
                    "series_count": result['series_count'],
                    "category_count": result['category_count'],
                },
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_list_images(self, params: ListImagesParams) -> ToolOutput:
        """列出演示文稿中的所有图像或特定幻灯片中的图像"""
        slide_filter = f"[{params.slide_index}]" if params.slide_index is not None else "None"

        script = f'''
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    prs = Presentation("{params.file_path}")
    images = []

    slide_indices = {slide_filter} if {slide_filter} != None else list(range(len(prs.slides)))

    for slide_idx in slide_indices:
        slide = prs.slides[slide_idx]
        img_idx = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({{
                    "slide_index": slide_idx,
                    "image_index": img_idx,
                    "width": shape.width / 914400,
                    "height": shape.height / 914400,
                    "left": shape.left / 914400,
                    "top": shape.top / 914400,
                }})
                img_idx += 1

    result = {{"success": True, "images": images}}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            images = result['images']
            summary = f"Found {len(images)} image(s) in {params.file_path}\n\n"
            for img in images:
                summary += f"Slide {img['slide_index']}, Image {img['image_index']}: "
                summary += f'{img["width"]:.2f}"x{img["height"]:.2f}" at ({img["left"]:.2f}", {img["top"]:.2f}")\n'

            return ToolOutput(
                blocks=[TextBlock(text=summary)],
                metadata={"file_path": params.file_path, "images": images},
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )

    @tool
    async def powerpoint_modify_image(self, params: ModifyImageParams) -> ToolOutput:
        """调整或替换幻灯片上的现有图像"""
        width_code = f"image_shape.width = int(Inches({params.width}))" if params.width else ""
        height_code = f"image_shape.height = int(Inches({params.height}))" if params.height else ""

        script = f'''
import json
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    prs = Presentation("{params.file_path}")
    slide = prs.slides[{params.slide_index}]

    image_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    if {params.image_index} < 0 or {params.image_index} >= len(image_shapes):
        raise ValueError(f"Image index {params.image_index} out of range")

    image_shape = image_shapes[{params.image_index}]

    {width_code}
    {height_code}

    new_width = image_shape.width / 914400
    new_height = image_shape.height / 914400

    prs.save("{params.file_path}")

    result = {{
        "success": True,
        "width": new_width,
        "height": new_height,
    }}
    print(json.dumps(result))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''

        cmd = f"python3 << 'PYSCRIPT_EOF'\n{script}\nPYSCRIPT_EOF"
        output, exit_code = await self.sandbox.run(cmd, max_bytes=1_000_000)
        try:
            result = json.loads(output)

            if not result.get("success"):
                return ToolOutput(
                    blocks=[TextBlock(text=f"❌ Error: {result.get('error')}")],
                    metadata={"error": result.get("error")},
                    reward=0.0,
                    finished=False,
                )

            return ToolOutput(
                blocks=[TextBlock(text=f"✅ Image modified on slide {params.slide_index}")],
                metadata={
                    "file_path": params.file_path,
                    "slide_index": params.slide_index,
                    "image_index": params.image_index,
                    "new_dimensions": {"width": result['width'], "height": result['height']},
                },
                reward=0.0,
                finished=False,
            )
        except json.JSONDecodeError:
            return ToolOutput(
                blocks=[TextBlock(text=f"❌ Error parsing output: {output}")],
                metadata={"error": "JSON decode failed", "output": output},
                reward=0.0,
                finished=False,
            )
